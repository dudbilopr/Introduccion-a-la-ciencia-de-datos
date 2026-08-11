import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_dark_theme(slide):
    # Set background color to dark gray
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 30, 35)

def format_text(shape, text, size_pt, bold=False, rgb=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = rgb
            run.font.name = 'Century Gothic'

def add_bullet(tf, text, level=0, bold=False, rgb=RGBColor(220, 220, 220), size=24):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb
    p.font.name = 'Century Gothic'

def main():
    prs = Presentation()
    # Cambiar a formato 16:9 (13.33 x 7.5 pulgadas)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Paths to generated images
    img_cover = r"C:\Users\dudbi\.gemini\antigravity-ide\brain\515cee07-6e58-4560-b901-e69205500e7c\ds_cover_1785838791043.png"
    img_github = r"C:\Users\dudbi\.gemini\antigravity-ide\brain\515cee07-6e58-4560-b901-e69205500e7c\ds_github_1785838801063.png"
    img_ai = r"C:\Users\dudbi\.gemini\antigravity-ide\brain\515cee07-6e58-4560-b901-e69205500e7c\ds_ai_1785838810770.png"

    # --- SLIDE 1: Title (High Impact) ---
    slide = prs.slides.add_slide(blank_layout)
    if os.path.exists(img_cover):
        # Insert background image full screen
        slide.shapes.add_picture(img_cover, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Title Box
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    format_text(tb1, "Introducción a la Ciencia de Datos", size_pt=60, bold=True, rgb=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    
    # Subtitle Box
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    format_text(tb2, "Clase 1: Reglas, Herramientas y Visión a Futuro\nDocente: Dudbil Olvasada Pabon Riaño", size_pt=30, bold=False, rgb=RGBColor(0, 204, 255), align=PP_ALIGN.CENTER)


    # --- Helper to create standard content slide ---
    def add_content_slide(title_text):
        sl = prs.slides.add_slide(blank_layout)
        apply_dark_theme(sl)
        # Title
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        format_text(tb, title_text, size_pt=45, bold=True, rgb=RGBColor(0, 255, 204))
        # Body frame
        body = sl.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        return sl, tf

    # --- SLIDE 2: Reglas ---
    sl, tf = add_content_slide("Reglas y Acuerdos de Clase")
    add_bullet(tf, "Para un ambiente de aprendizaje seguro y productivo:", bold=True, rgb=RGBColor(255,255,255), size=28)
    add_bullet(tf, "Puntualidad: Iniciamos a las 7:00 AM en punto.", level=1)
    add_bullet(tf, "Descansos Cognitivos: Haremos una pausa para evitar fatiga.", level=1)
    add_bullet(tf, "Estrategia EMI: Tendremos inmersiones de 10 min en inglés.", level=1)
    add_bullet(tf, "Respeto: Ambiente seguro para compartir ideas y debatir.", level=1)
    add_bullet(tf, "Comunidad: Utilizaremos un Espacio en Gmail para dudas.", level=1)

    # --- SLIDE 3: EMI Warm up ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Warm-Up")
    add_bullet(tf, "Let's activate our English!", bold=True, rgb=RGBColor(255,255,255), size=36)
    add_bullet(tf, "\nQuestion:", level=0, size=32)
    add_bullet(tf, "What is the first word that comes to your mind when you hear 'Data Science'?", level=1, bold=True, rgb=RGBColor(0, 204, 255), size=40)

    # --- SLIDE 4: EMI Video ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Video Activity")
    add_bullet(tf, "Objectives:", bold=True, rgb=RGBColor(255,255,255), size=28)
    add_bullet(tf, "1. Learn key technical vocabulary.", level=1)
    add_bullet(tf, "2. Understand the core concept of Data Science.", level=1)
    add_bullet(tf, "\nVideo: What is Data Science?", bold=True, rgb=RGBColor(0, 255, 204), size=28)
    add_bullet(tf, "Link: https://www.youtube.com/watch?v=xN6dbHSbohk", level=1)

    # --- SLIDE 5: EMI Vocab ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Vocabulary Guide")
    add_bullet(tf, "Key terms from the video:", bold=True, rgb=RGBColor(255,255,255), size=28)
    add_bullet(tf, "Data: Raw facts and statistics.", level=1)
    add_bullet(tf, "Insights: Deep understanding or discoveries.", level=1)
    add_bullet(tf, "Machine Learning: Algorithms that learn from data.", level=1)
    add_bullet(tf, "Predict: To say what will happen in the future.", level=1)

    # --- SLIDE 6: Mentimeter 1 ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Let's Check! (1/2)")
    add_bullet(tf, "Join the Mentimeter! (Questions 1-3)", bold=True, rgb=RGBColor(255,255,255), size=28)
    
    add_bullet(tf, "Q1: What are the three V's of Big Data mentioned in the video?", level=0, size=24, bold=True)
    add_bullet(tf, "A) Video, Vision, Value | B) Volume, Velocity, Variety | C) Virtual, Vector, Volume", level=1, size=20, rgb=RGBColor(0,204,255))
    
    add_bullet(tf, "Q2: What is the first component of Data Science according to the video?", level=0, size=24, bold=True)
    add_bullet(tf, "A) Domain expertise and scientific methods | B) Buying hardware | C) Writing blog posts", level=1, size=20, rgb=RGBColor(0,204,255))

    add_bullet(tf, "Q3: What does 'Velocity' refer to in Big Data?", level=0, size=24, bold=True)
    add_bullet(tf, "A) The speed of the computer | B) Huge amounts of data flowing at a tremendous speed | C) Typing skills", level=1, size=20, rgb=RGBColor(0,204,255))

    # --- SLIDE 7: Mentimeter 2 ---
    sl, tf = add_content_slide("🇬🇧 EMI Strategy: Let's Check! (2/2)")
    add_bullet(tf, "Join the Mentimeter! (Questions 4-5)", bold=True, rgb=RGBColor(255,255,255), size=28)
    
    add_bullet(tf, "Q4: What helps provide insights by building and running automated models?", level=0, size=24, bold=True)
    add_bullet(tf, "A) Machine Learning | B) Data Cleansing | C) Web Design", level=1, size=20, rgb=RGBColor(0,204,255))

    add_bullet(tf, "Q5: (Word Cloud) Write one format of Big Data mentioned in the video!", level=0, size=24, bold=True)
    add_bullet(tf, "Example: Structured, Semi-structured, Unstructured, JSON, XML...", level=1, size=20, rgb=RGBColor(0,204,255))

    # --- SLIDE 8: GitHub ---
    sl = prs.slides.add_slide(blank_layout)
    apply_dark_theme(sl)
    if os.path.exists(img_github):
        sl.shapes.add_picture(img_github, Inches(8.33), Inches(1.5), width=Inches(4.5), height=Inches(4.5))
    
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    format_text(tb, "Herramientas Tecnológicas: GitHub", size_pt=45, bold=True, rgb=RGBColor(0, 255, 204))
    
    body = sl.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7.5), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    add_bullet(tf, "¿Por qué GitHub?", bold=True, rgb=RGBColor(255,255,255), size=32)
    add_bullet(tf, "Es el portafolio profesional estándar para código.", level=1, size=24)
    add_bullet(tf, "Paso a paso:", bold=True, rgb=RGBColor(255,255,255), size=32)
    add_bullet(tf, "1. Ingresa a github.com.", level=1, size=24)
    add_bullet(tf, "2. Haz clic en 'Sign Up' y usa tu correo institucional.", level=1, size=24)
    add_bullet(tf, "3. Personaliza tu perfil (Foto, Bio).", level=1, size=24)

    # --- SLIDE 9: LLMNotebook ---
    sl = prs.slides.add_slide(blank_layout)
    apply_dark_theme(sl)
    if os.path.exists(img_ai):
        sl.shapes.add_picture(img_ai, Inches(8.33), Inches(1.5), width=Inches(4.5), height=Inches(4.5))
    
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    format_text(tb, "Herramientas Tecnológicas: LLMNotebook", size_pt=45, bold=True, rgb=RGBColor(0, 255, 204))
    
    body = sl.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7.5), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    add_bullet(tf, "Asistencia con Inteligencia Artificial", bold=True, rgb=RGBColor(255,255,255), size=32)
    add_bullet(tf, "Usaremos IA para acelerar el aprendizaje, no para reemplazarlo.", level=1, size=24)
    add_bullet(tf, "Paso a paso:", bold=True, rgb=RGBColor(255,255,255), size=32)
    add_bullet(tf, "1. Accede a la plataforma indicada.", level=1, size=24)
    add_bullet(tf, "2. Formula un 'prompt' efectivo (Contexto + Tarea + Formato).", level=1, size=24)
    add_bullet(tf, "3. Analiza críticamente; tú eres el piloto, la IA es el copiloto.", level=1, size=24)

    # --- SLIDE 10: Actividad Grupal ---
    sl, tf = add_content_slide("Actividad Central: El Futuro del Científico")
    add_bullet(tf, "Objetivo: Analizar reportes clave (Donoho, WEF, Stanford AI).", bold=True, rgb=RGBColor(255,255,255), size=28)
    add_bullet(tf, "Dinámica en grupos:", size=24, bold=True, rgb=RGBColor(0, 204, 255))
    add_bullet(tf, "Distribuyan los documentos y usen LLMNotebook para sintetizar información.", level=1)
    add_bullet(tf, "Entregables de la exposición:", size=24, bold=True, rgb=RGBColor(0, 204, 255))
    add_bullet(tf, "1. Resumen visual: Diapositiva/infografía con puntos clave.", level=1)
    add_bullet(tf, "2. Conexión práctica: Un ejemplo real del mercado laboral.", level=1)
    add_bullet(tf, "3. Dato de Impacto: Algo que los haya sorprendido genuinamente.", level=1)

    # --- SLIDE 11: Tiempos Exposición ---
    sl, tf = add_content_slide("Sustentaciones: Reglas de Juego")
    add_bullet(tf, "Lineamientos:", bold=True, rgb=RGBColor(255,255,255), size=32)
    add_bullet(tf, "Tiempo: 15 minutos exactos por grupo.", level=1, size=28)
    add_bullet(tf, "Dinámica: Participación equitativa y dominio del tema.", level=1, size=28)
    add_bullet(tf, "Feedback: Preguntas cruzadas y co-evaluación al final.", level=1, size=28)
    add_bullet(tf, "\n¡Manos a la obra!", bold=True, rgb=RGBColor(0, 255, 204), size=40, level=0)

    prs.save('Presentacion_Premium_Clase1.pptx')
    print("Presentación Premium guardada con éxito.")

if __name__ == '__main__':
    main()
