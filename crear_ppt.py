from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def main():
    prs = Presentation()
    
    # 0 = Title slide
    # 1 = Title and Content
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Introducción a la Ciencia de Datos"
    subtitle.text = "Clase 1: Reglas, Herramientas y Visión a Futuro\nDocente: Dudbil Olvasada Pabon Riaño"

    # Slide 2: Reglas de Clase
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Reglas y Acuerdos de Clase"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Para un ambiente de aprendizaje seguro y productivo:"
    p = tf.add_paragraph()
    p.text = "Puntualidad: Iniciamos a las 7:00 AM en punto."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Descansos Cognitivos: Haremos una pausa para evitar fatiga."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Estrategia EMI: Tendremos inmersiones de 10 min en inglés."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Respeto: Ambiente seguro para compartir ideas y debatir."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Canal de Comunicación: Utilizaremos un Espacio en Gmail."
    p.level = 1

    # Slide 3: EMI - Warm-up
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "🇬🇧 EMI Strategy: Warm-Up"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Let's activate our English!"
    p = tf.add_paragraph()
    p.text = "Question:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "What is the first word that comes to your mind when you hear 'Data Science'?"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.level = 2

    # Slide 4: EMI - Objectives & Video
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "🇬🇧 EMI Strategy: Video Activity"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Objectives:"
    p = tf.add_paragraph()
    p.text = "1. Learn key technical vocabulary."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Understand the core concept of Data Science."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Video: What is Data Science?"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Link: https://www.youtube.com/watch?v=xN6dbHSbohk"
    p.level = 1

    # Slide 5: EMI - Vocabulary Guide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "🇬🇧 EMI Strategy: Vocabulary Guide"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Key terms from the video:"
    p = tf.add_paragraph()
    p.text = "Data: Raw facts and statistics."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Insights: Deep understanding or discoveries."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Machine Learning: Algorithms that learn from data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Predict: To say what will happen in the future."
    p.level = 1

    # Slide 6: EMI - Mentimeter (Part 1)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "🇬🇧 EMI Strategy: Let's Check! (1/2)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Join the Mentimeter! (Questions 1-3)"
    
    p = tf.add_paragraph()
    p.text = "Q1: What are the three V's of Big Data mentioned in the video?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A) Video, Vision, Value | B) Volume, Velocity, Variety | C) Virtual, Vector, Volume"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Q2: What is the first component of Data Science according to the video?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A) Domain expertise and scientific methods | B) Buying hardware | C) Writing blog posts"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Q3: What does 'Velocity' refer to in Big Data?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A) The speed of the computer | B) Huge amounts of data flowing at a tremendous speed | C) Typing skills"
    p.level = 2

    # Slide 7: EMI - Mentimeter (Part 2)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "🇬🇧 EMI Strategy: Let's Check! (2/2)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Join the Mentimeter! (Questions 4-5)"
    
    p = tf.add_paragraph()
    p.text = "Q4: What helps provide insights by building and running automated models?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "A) Machine Learning | B) Data Cleansing | C) Web Design"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Q5: (Word Cloud) Write one format of Big Data mentioned in the video!"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Example: Structured, Semi-structured, Unstructured, JSON, XML..."
    p.level = 2

    # Slide 7: Herramientas - GitHub
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Herramientas Tecnológicas: GitHub"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "¿Por qué GitHub?"
    p = tf.add_paragraph()
    p.text = "Es el portafolio profesional estándar para código y proyectos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Paso a paso:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "1. Ingresa a github.com."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Haz clic en 'Sign Up' y regístrate con tu correo."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Personaliza tu perfil."
    p.level = 1

    # Slide 8: Herramientas - LLMNotebook
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Herramientas Tecnológicas: LLMNotebook"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Asistencia con Inteligencia Artificial"
    p = tf.add_paragraph()
    p.text = "Usaremos LLMNotebook para acelerar nuestro aprendizaje."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Paso a paso:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "1. Accede a la plataforma proporcionada por el docente."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Formula un 'prompt' claro. (Ej: 'Actúa como Data Scientist y resume en 3 puntos...') "
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Analiza críticamente la respuesta; la IA es tu asistente, tú tomas las decisiones."
    p.level = 1

    # Slide 9: Actividad Grupal
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Actividad Central: El Futuro del Científico de Datos"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Objetivo: Analizar reportes clave (Donoho, WEF, Stanford AI)."
    p = tf.add_paragraph()
    p.text = "Trabajo en grupos:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Distribúyanse los documentos y usen LLMNotebook para procesar la información rápido."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Entregables de la exposición:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "1. Resumen visual: Una diapositiva o infografía con los puntos clave."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Conexión práctica: Un ejemplo real aplicable."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Dato de Impacto: Un hallazgo sorprendente."
    p.level = 1

    # Slide 10: Tiempos Exposición
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Sustentaciones"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Reglas de la Exposición:"
    p = tf.add_paragraph()
    p.text = "Tiempo: 15 minutos exactos por grupo."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dinámica: Todos deben participar y mostrar dominio del tema."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Feedback: Al final habrá espacio para preguntas y co-evaluación formativa."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "¡Manos a la obra!"
    p.level = 0

    prs.save('Presentacion_Clase1.pptx')
    print("Presentación guardada con éxito.")

if __name__ == '__main__':
    main()
